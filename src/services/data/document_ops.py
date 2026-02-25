"""
src/services/data/document_ops.py
===================================
Complete document processing for Excel, Word, PDF, PowerPoint.
- openpyxl: Excel file-level read/write (no Excel needed)
- xlwings: Live Excel automation (formulas, pivots)
- python-docx: Word documents (read, write, format, templates)
- PyMuPDF (fitz): PDF extraction (10x faster than PyPDF2)
- python-pptx: PowerPoint creation and editing
"""

import os
import re
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path

# Excel
import openpyxl
from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
from openpyxl.utils import get_column_letter

# Excel Live (optional - requires Excel installation)
try:
    import xlwings as xw
    XLWINGS_AVAILABLE = True
except ImportError:
    XLWINGS_AVAILABLE = False
    print("[DocOps] xlwings not available - live Excel features disabled")

# Word
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH

# PDF
import fitz  # PyMuPDF

# PowerPoint
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt


# ─────────────────────────────────────────────────────
# EXCEL OPERATIONS (openpyxl - No Excel needed)
# ─────────────────────────────────────────────────────

class ExcelOps:
    """Excel operations using openpyxl (pure Python, no Excel needed)."""
    
    @staticmethod
    def read_data(
        file_path: str,
        sheet_name: Optional[str] = None,
        max_rows: int = 100
    ) -> List[List[Any]]:
        """Reads Excel data as 2D list."""
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            ws = wb[sheet_name] if sheet_name else wb.active
            
            data = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= max_rows:
                    break
                data.append(list(row))
            
            wb.close()
            return data
        except Exception as e:
            print(f"[Excel] Read error: {e}")
            return []
    
    @staticmethod
    def write_data(
        file_path: str,
        data: List[List[Any]],
        sheet_name: str = "Sheet1"
    ) -> str:
        """Writes 2D list to Excel."""
        try:
            # Create or load workbook
            if os.path.exists(file_path):
                wb = openpyxl.load_workbook(file_path)
            else:
                wb = openpyxl.Workbook()
            
            # Get or create sheet
            if sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
            else:
                ws = wb.create_sheet(sheet_name)
            
            # Write data
            for row_idx, row_data in enumerate(data, start=1):
                for col_idx, value in enumerate(row_data, start=1):
                    ws.cell(row=row_idx, column=col_idx, value=value)
            
            wb.save(file_path)
            wb.close()
            return f"✅ Wrote {len(data)} rows to '{file_path}'"
        except Exception as e:
            return f"❌ Write failed: {e}"
    
    @staticmethod
    def update_cell(
        file_path: str,
        cell: str,
        value: Any,
        sheet_name: Optional[str] = None
    ) -> str:
        """Updates single cell (e.g., A1, B5)."""
        try:
            wb = openpyxl.load_workbook(file_path)
            ws = wb[sheet_name] if sheet_name else wb.active
            ws[cell] = value
            wb.save(file_path)
            wb.close()
            return f"✅ Updated {cell} = {value}"
        except Exception as e:
            return f"❌ Update failed: {e}"
    
    @staticmethod
    def style_cells(
        file_path: str,
        cell_range: str,
        bg_color: Optional[str] = None,
        font_color: Optional[str] = None,
        bold: bool = False,
        font_size: Optional[int] = None,
        border: bool = False,
        sheet_name: Optional[str] = None
    ) -> str:
        """Applies styling to cell range."""
        try:
            wb = openpyxl.load_workbook(file_path)
            ws = wb[sheet_name] if sheet_name else wb.active
            
            cells_styled = 0
            for row in ws[cell_range]:
                for cell in row:
                    # Background color
                    if bg_color:
                        cell.fill = PatternFill(
                            start_color=bg_color.replace('#', ''),
                            end_color=bg_color.replace('#', ''),
                            fill_type="solid"
                        )
                    
                    # Font
                    font_kwargs = {}
                    if font_color:
                        font_kwargs['color'] = font_color.replace('#', '')
                    if bold:
                        font_kwargs['bold'] = True
                    if font_size:
                        font_kwargs['size'] = font_size
                    
                    if font_kwargs:
                        cell.font = Font(**font_kwargs)
                    
                    # Border
                    if border:
                        thin = Side(style='thin', color='000000')
                        cell.border = Border(
                            left=thin, right=thin, top=thin, bottom=thin
                        )
                    
                    cells_styled += 1
            
            wb.save(file_path)
            wb.close()
            return f"✅ Styled {cells_styled} cells"
        except Exception as e:
            return f"❌ Styling failed: {e}"
    
    @staticmethod
    def create_table(
        file_path: str,
        headers: List[str],
        data: List[List[Any]],
        table_name: str = "MyTable",
        sheet_name: str = "Sheet1"
    ) -> str:
        """Creates formatted table with headers."""
        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = sheet_name
            
            # Write headers with styling
            for col_idx, header in enumerate(headers, start=1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = Font(bold=True, color='FFFFFF')
                cell.fill = PatternFill(start_color='4472C4', end_color='4472C4', fill_type='solid')
                cell.alignment = Alignment(horizontal='center')
            
            # Write data
            for row_idx, row_data in enumerate(data, start=2):
                for col_idx, value in enumerate(row_data, start=1):
                    ws.cell(row=row_idx, column=col_idx, value=value)
            
            # Auto-size columns
            for column in ws.columns:
                max_length = 0
                column_letter = get_column_letter(column[0].column)
                for cell in column:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                ws.column_dimensions[column_letter].width = min(max_length + 2, 50)
            
            wb.save(file_path)
            wb.close()
            return f"✅ Created table with {len(data)} rows"
        except Exception as e:
            return f"❌ Table creation failed: {e}"


# ─────────────────────────────────────────────────────
# EXCEL LIVE (xlwings - Requires Excel installation)
# ─────────────────────────────────────────────────────

class ExcelLive:
    """Live Excel automation with formula recalculation and pivot refresh."""
    
    @staticmethod
    def refresh_formulas(file_path: str) -> str:
        """Forces formula recalculation."""
        if not XLWINGS_AVAILABLE:
            return "❌ xlwings not available"
        
        try:
            app = xw.App(visible=False)
            wb = app.books.open(file_path)
            wb.app.calculate()  # Force calculate
            wb.save()
            wb.close()
            app.quit()
            return f"✅ Formulas recalculated in '{file_path}'"
        except Exception as e:
            return f"❌ Formula refresh failed: {e}"
    
    @staticmethod
    def refresh_pivot_tables(file_path: str) -> str:
        """Refreshes all pivot tables."""
        if not XLWINGS_AVAILABLE:
            return "❌ xlwings not available"
        
        try:
            app = xw.App(visible=False)
            wb = app.books.open(file_path)
            
            pivot_count = 0
            for sheet in wb.sheets:
                for pivot in sheet.api.PivotTables():
                    pivot.PivotCache().Refresh()
                    pivot_count += 1
            
            wb.save()
            wb.close()
            app.quit()
            return f"✅ Refreshed {pivot_count} pivot table(s)"
        except Exception as e:
            return f"❌ Pivot refresh failed: {e}"
    
    @staticmethod
    def run_macro(file_path: str, macro_name: str) -> str:
        """Runs VBA macro."""
        if not XLWINGS_AVAILABLE:
            return "❌ xlwings not available"
        
        try:
            app = xw.App(visible=False)
            wb = app.books.open(file_path)
            app.api.Run(macro_name)
            wb.save()
            wb.close()
            app.quit()
            return f"✅ Ran macro: {macro_name}"
        except Exception as e:
            return f"❌ Macro execution failed: {e}"


# ─────────────────────────────────────────────────────
# WORD OPERATIONS (python-docx)
# ─────────────────────────────────────────────────────

class WordOps:
    """Word document operations."""
    
    @staticmethod
    def read_text(file_path: str) -> str:
        """Extracts all text from Word document."""
        try:
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text
        except Exception as e:
            return f"Error: {e}"
    
    @staticmethod
    def extract_with_pattern(file_path: str, pattern: str) -> List[str]:
        """Extracts text matching regex pattern."""
        try:
            doc = Document(file_path)
            full_text = "\n".join([para.text for para in doc.paragraphs])
            matches = re.findall(pattern, full_text, re.IGNORECASE)
            return matches
        except Exception as e:
            print(f"[Word] Extract error: {e}")
            return []
    
    @staticmethod
    def fill_template(
        template_path: str,
        output_path: str,
        data: Dict[str, str]
    ) -> str:
        """
        Fills Word template by replacing {{placeholders}}.
        
        Example:
            template: "Dear {{name}}, your total is {{amount}}."
            data: {"name": "John", "amount": "$100"}
            result: "Dear John, your total is $100."
        """
        try:
            doc = Document(template_path)
            
            replacements = 0
            
            # Replace in paragraphs
            for para in doc.paragraphs:
                for key, value in data.items():
                    for fmt in [f"{{{{{key}}}}}", f"{{{key}}}"]:
                        if fmt in para.text:
                            para.text = para.text.replace(fmt, str(value))
                            replacements += 1
            
            # Replace in tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for key, value in data.items():
                            for fmt in [f"{{{{{key}}}}}", f"{{{key}}}"]:
                                if fmt in cell.text:
                                    cell.text = cell.text.replace(fmt, str(value))
                                    replacements += 1
            
            doc.save(output_path)
            return f"✅ Filled template: {replacements} replacements made"
        except Exception as e:
            return f"❌ Template fill failed: {e}"
    
    @staticmethod
    def format_text(
        file_path: str,
        target_text: str,
        font_name: Optional[str] = None,
        font_size: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        color: Optional[str] = None
    ) -> str:
        """Applies formatting to text."""
        try:
            doc = Document(file_path)
            changes = 0
            
            color_map = {
                "red": (255, 0, 0),
                "blue": (0, 0, 255),
                "green": (0, 128, 0),
                "black": (0, 0, 0),
            }
            
            for para in doc.paragraphs:
                if target_text.lower() in para.text.lower():
                    for run in para.runs:
                        if font_name:
                            run.font.name = font_name
                        if font_size:
                            run.font.size = Pt(font_size)
                        if bold is not None:
                            run.font.bold = bold
                        if italic is not None:
                            run.font.italic = italic
                        if color and color.lower() in color_map:
                            run.font.color.rgb = RGBColor(*color_map[color.lower()])
                        changes += 1
            
            doc.save(file_path)
            return f"✅ Formatted {changes} text runs"
        except Exception as e:
            return f"❌ Formatting failed: {e}"
    
    @staticmethod
    def create_document(
        output_path: str,
        title: str,
        sections: List[Dict[str, str]]
    ) -> str:
        """
        Creates new Word document with structure.
        
        sections: [
            {"heading": "Introduction", "text": "This is..."},
            {"heading": "Results", "text": "We found..."}
        ]
        """
        try:
            doc = Document()
            
            # Add title
            title_para = doc.add_heading(title, level=0)
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Add sections
            for section in sections:
                if "heading" in section:
                    doc.add_heading(section["heading"], level=1)
                if "text" in section:
                    doc.add_paragraph(section["text"])
            
            doc.save(output_path)
            return f"✅ Created document: {output_path}"
        except Exception as e:
            return f"❌ Document creation failed: {e}"


# ─────────────────────────────────────────────────────
# PDF OPERATIONS (PyMuPDF/fitz - 10x faster)
# ─────────────────────────────────────────────────────

class PDFOps:
    """PDF operations using PyMuPDF (fitz) - much faster than PyPDF2."""
    
    @staticmethod
    def extract_text(file_path: str, max_pages: int = 50) -> str:
        """Extracts all text from PDF."""
        try:
            doc = fitz.open(file_path)
            text = ""
            for page_num in range(min(len(doc), max_pages)):
                page = doc[page_num]
                text += page.get_text() + "\n"
            doc.close()
            return text
        except Exception as e:
            return f"Error: {e}"
    
    @staticmethod
    def extract_with_pattern(file_path: str, pattern: str) -> List[str]:
        """Extracts text matching regex pattern."""
        try:
            text = PDFOps.extract_text(file_path)
            matches = re.findall(pattern, text, re.IGNORECASE)
            return matches
        except Exception as e:
            print(f"[PDF] Extract error: {e}")
            return []
    
    @staticmethod
    def extract_tables(file_path: str) -> List[List[List[str]]]:
        """
        Extracts tables from PDF.
        Returns: [page1_tables, page2_tables, ...]
        """
        try:
            doc = fitz.open(file_path)
            all_tables = []
            
            for page in doc:
                tables = page.find_tables()
                page_tables = []
                for table in tables:
                    # Extract table as 2D list
                    table_data = table.extract()
                    page_tables.append(table_data)
                all_tables.append(page_tables)
            
            doc.close()
            return all_tables
        except Exception as e:
            print(f"[PDF] Table extraction error: {e}")
            return []
    
    @staticmethod
    def convert_to_images(
        file_path: str,
        output_dir: str,
        dpi: int = 150
    ) -> str:
        """Converts PDF pages to images."""
        try:
            doc = fitz.open(file_path)
            os.makedirs(output_dir, exist_ok=True)
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                pix = page.get_pixmap(dpi=dpi)
                output_path = os.path.join(output_dir, f"page_{page_num + 1}.png")
                pix.save(output_path)
            
            doc.close()
            return f"✅ Converted {len(doc)} pages to images in '{output_dir}'"
        except Exception as e:
            return f"❌ Conversion failed: {e}"
    
    @staticmethod
    def merge_pdfs(input_paths: List[str], output_path: str) -> str:
        """Merges multiple PDFs."""
        try:
            result = fitz.open()
            
            for pdf_path in input_paths:
                pdf = fitz.open(pdf_path)
                result.insert_pdf(pdf)
                pdf.close()
            
            result.save(output_path)
            result.close()
            return f"✅ Merged {len(input_paths)} PDFs into '{output_path}'"
        except Exception as e:
            return f"❌ Merge failed: {e}"
    
    @staticmethod
    def split_pdf(
        input_path: str,
        output_dir: str,
        pages_per_file: int = 1
    ) -> str:
        """Splits PDF into multiple files."""
        try:
            doc = fitz.open(input_path)
            os.makedirs(output_dir, exist_ok=True)
            
            file_num = 1
            for start_page in range(0, len(doc), pages_per_file):
                end_page = min(start_page + pages_per_file, len(doc))
                
                output_pdf = fitz.open()
                output_pdf.insert_pdf(doc, from_page=start_page, to_page=end_page - 1)
                
                output_path = os.path.join(output_dir, f"part_{file_num}.pdf")
                output_pdf.save(output_path)
                output_pdf.close()
                file_num += 1
            
            doc.close()
            return f"✅ Split into {file_num - 1} files in '{output_dir}'"
        except Exception as e:
            return f"❌ Split failed: {e}"


# ─────────────────────────────────────────────────────
# POWERPOINT OPERATIONS (python-pptx)
# ─────────────────────────────────────────────────────

class PowerPointOps:
    """PowerPoint operations."""
    
    @staticmethod
    def read_text(file_path: str) -> str:
        """Extracts all text from PowerPoint."""
        try:
            prs = Presentation(file_path)
            text = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            
            return "\n".join(text)
        except Exception as e:
            return f"Error: {e}"
    
    @staticmethod
    def create_presentation(
        output_path: str,
        title: str,
        slides: List[Dict[str, Any]]
    ) -> str:
        """
        Creates PowerPoint presentation.
        
        slides: [
            {"title": "Slide 1", "content": "Bullet 1\nBullet 2"},
            {"title": "Slide 2", "content": "More content"}
        ]
        """
        try:
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            
            # Content slides
            for slide_data in slides:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                if "title" in slide_data:
                    slide.shapes.title.text = slide_data["title"]
                
                if "content" in slide_data:
                    text_frame = slide.placeholders[1].text_frame
                    for line in slide_data["content"].split("\n"):
                        p = text_frame.add_paragraph()
                        p.text = line
                        p.level = 0
            
            prs.save(output_path)
            return f"✅ Created presentation with {len(slides) + 1} slides"
        except Exception as e:
            return f"❌ Presentation creation failed: {e}"
    
    @staticmethod
    def add_slide_with_image(
        file_path: str,
        title: str,
        image_path: str
    ) -> str:
        """Adds slide with title and image."""
        try:
            if os.path.exists(file_path):
                prs = Presentation(file_path)
            else:
                prs = Presentation()
            
            blank_slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add title
            left = PptxInches(0.5)
            top = PptxInches(0.5)
            width = PptxInches(9)
            height = PptxInches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = title
            
            # Add image
            left = PptxInches(1)
            top = PptxInches(2)
            slide.shapes.add_picture(image_path, left, top, height=PptxInches(4))
            
            prs.save(file_path)
            return f"✅ Added slide with image"
        except Exception as e:
            return f"❌ Slide addition failed: {e}"


# ─────────────────────────────────────────────────────
# FORMAT CONVERSION
# ─────────────────────────────────────────────────────

class FormatConverter:
    """Converts between document formats."""
    
    @staticmethod
    def excel_to_pdf(excel_path: str, pdf_path: str) -> str:
        """Converts Excel to PDF (requires xlwings)."""
        if not XLWINGS_AVAILABLE:
            return "❌ xlwings not available - Excel to PDF requires Excel installation"
        
        try:
            app = xw.App(visible=False)
            wb = app.books.open(excel_path)
            wb.to_pdf(pdf_path)
            wb.close()
            app.quit()
            return f"✅ Converted to PDF: {pdf_path}"
        except Exception as e:
            return f"❌ Conversion failed: {e}"
    
    @staticmethod
    def csv_to_excel(csv_path: str, excel_path: str) -> str:
        """Converts CSV to Excel."""
        try:
            import csv
            
            wb = openpyxl.Workbook()
            ws = wb.active
            
            with open(csv_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                for row in reader:
                    ws.append(row)
            
            wb.save(excel_path)
            wb.close()
            return f"✅ Converted CSV to Excel: {excel_path}"
        except Exception as e:
            return f"❌ Conversion failed: {e}"
    
    @staticmethod
    def excel_to_csv(excel_path: str, csv_path: str, sheet_name: Optional[str] = None) -> str:
        """Converts Excel to CSV."""
        try:
            import csv
            
            wb = openpyxl.load_workbook(excel_path, data_only=True)
            ws = wb[sheet_name] if sheet_name else wb.active
            
            with open(csv_path, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                for row in ws.iter_rows(values_only=True):
                    writer.writerow(row)
            
            wb.close()
            return f"✅ Converted Excel to CSV: {csv_path}"
        except Exception as e:
            return f"❌ Conversion failed: {e}"


# ─────────────────────────────────────────────────────
# GLOBAL INSTANCES
# ─────────────────────────────────────────────────────

excel_ops = ExcelOps()
excel_live = ExcelLive()
word_ops = WordOps()
pdf_ops = PDFOps()
ppt_ops = PowerPointOps()
format_converter = FormatConverter()
